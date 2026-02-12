#!/usr/bin/env python3
"""Generate OCE Active Tickets PowerPoint Presentation for Blend Engineering."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
import datetime

# ── Colors ──────────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x07, 0x0B, 0x14)
BG_CARD      = RGBColor(0x11, 0x18, 0x27)
TEXT_WHITE    = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_GRAY     = RGBColor(0x94, 0xA3, 0xB8)
TEXT_DIM      = RGBColor(0x64, 0x74, 0x8B)
RED           = RGBColor(0xF8, 0x71, 0x71)
ORANGE        = RGBColor(0xFB, 0x92, 0x3C)
AMBER         = RGBColor(0xFB, 0xBF, 0x24)
GREEN         = RGBColor(0x34, 0xD3, 0x99)
BLUE          = RGBColor(0x60, 0xA5, 0xFA)
PURPLE        = RGBColor(0xA7, 0x8B, 0xFA)
SLATE         = RGBColor(0x94, 0xA3, 0xB8)
ACCENT_RED    = RGBColor(0xEF, 0x44, 0x44)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ─────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_kpi_card(slide, left, top, value, label, accent_color):
    card = add_shape(slide, left, top, Inches(1.8), Inches(1.2), BG_CARD, RGBColor(0x1E, 0x29, 0x3B))
    add_text_box(slide, left + Inches(0.1), top + Inches(0.15), Inches(1.6), Inches(0.6),
                 str(value), font_size=28, color=accent_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.72), Inches(1.6), Inches(0.4),
                 label, font_size=8, color=TEXT_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

def add_table_slide(slide, headers, rows, left, top, width, col_widths, row_height=Inches(0.35)):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, Inches(0.35) * (len(rows) + 1))
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = TEXT_GRAY
            paragraph.font.name = 'Segoe UI'
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0D, 0x11, 0x1E)

    # Rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.color.rgb = TEXT_WHITE
                paragraph.font.name = 'Segoe UI'
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD if r_idx % 2 == 0 else RGBColor(0x0D, 0x11, 0x1E)

    return table_shape

# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════════════════

slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
set_slide_bg(slide1, BG_DARK)

# Decorative top bar
add_shape(slide1, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_RED)

# Title
add_text_box(slide1, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
             'OCE Active Tickets Report', font_size=44, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide1, Inches(1), Inches(3.2), Inches(11.3), Inches(0.6),
             'On-Call Engineering  |  EPD-Logged  |  Production Issues Tracker', font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Projects line
add_text_box(slide1, Inches(1), Inches(4.0), Inches(11.3), Inches(0.5),
             'DD  ·  DDINDIA  ·  SENG  ·  CLN  ·  CBP  ·  DA  ·  RHL  ·  IMB  ·  APEX  ·  BAI', font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Blend Engineering label
add_text_box(slide1, Inches(1), Inches(5.2), Inches(11.3), Inches(0.4),
             'Blend Engineering', font_size=16, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# Date
now = datetime.datetime.now().strftime('%B %d, %Y')
add_text_box(slide1, Inches(1), Inches(5.7), Inches(11.3), Inches(0.4),
             f'Data sourced from JIRA (blendlabs.atlassian.net)  ·  {now}', font_size=11, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide1, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), ACCENT_RED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — Executive Summary (KPIs)
# ═══════════════════════════════════════════════════════════════════════════

slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)
add_shape(slide2, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_RED)

add_text_box(slide2, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             'Executive Summary — OCE Active Tickets', font_size=28, color=RED, bold=True)
add_text_box(slide2, Inches(0.5), Inches(0.9), Inches(12), Inches(0.3),
             'Key Performance Indicators across all projects', font_size=12, color=TEXT_GRAY)

# KPI cards - row 1
kpi_data = [
    ('131', 'Total Active\nOCE Tickets', RED),
    ('42', 'P0 / P1\nCritical', ORANGE),
    ('52', 'EPD-Logged\nIssues', AMBER),
    ('10', 'Projects\nAffected', BLUE),
    ('11', 'In Progress\n/ QA', PURPLE),
    ('7', 'Ready\nfor QA', GREEN),
]

start_x = Inches(0.7)
for i, (val, lbl, clr) in enumerate(kpi_data):
    x = start_x + Inches(i * 2.05)
    add_kpi_card(slide2, x, Inches(1.6), val, lbl, clr)

# Summary bullets
bullets = [
    ('Critical Finding:', 'CBP (35+) and DD (30) lead with most active OCE tickets. DD was previously uncounted. 5 CBP P1 issues unresolved including stuck CP enablement nodes.', RED),
    ('EPD Trends:', '52 EPD-logged production issues across 10 projects — DD (30), SENG (20), CLN (18), DDINDIA (8), DA/RHL (13), CBP (35+). DD is 2nd largest.', ORANGE),
    ('Staffing Risk:', 'Jacob Powers (DD: 10), Omar Aparicio (SENG: 12), Ya He (DD: 8) carry heavy OCE loads. DDINDIA: Gaurav Kumar & Kumar Utkarsh handle most tickets.', AMBER),
    ('Blocked Items:', '5 tickets blocked — CBP-11817, CBP-12877, CLN-185, RHL-1659, DDINDIA-555 (4506C document issue for Wyndham) — need escalation.', RED),
    ('Resolution Gap:', '77 of 131 tickets (59%) in "To Do/Ready for Dev" with no active development, indicating a significant resolution backlog across all teams.', BLUE),
]

y = Inches(3.2)
for title, desc, clr in bullets:
    # Accent bar
    add_shape(slide2, Inches(0.7), y + Inches(0.04), Inches(0.06), Inches(0.4), clr)
    add_text_box(slide2, Inches(0.95), y, Inches(2.0), Inches(0.5),
                 title, font_size=11, color=clr, bold=True)
    add_text_box(slide2, Inches(2.9), y, Inches(9.5), Inches(0.5),
                 desc, font_size=10, color=TEXT_WHITE)
    y += Inches(0.62)

# Bottom bar
add_shape(slide2, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), ACCENT_RED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — Active Tickets by Project
# ═══════════════════════════════════════════════════════════════════════════

slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_DARK)
add_shape(slide3, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_RED)

add_text_box(slide3, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             'Active Tickets by Project', font_size=28, color=RED, bold=True)

# Chart data
chart_data = CategoryChartData()
chart_data.categories = ['CBP', 'DD', 'SENG', 'CLN', 'DDINDIA', 'RHL', 'DA', 'APEX', 'BAI', 'IMB']
chart_data.add_series('Active Tickets', (35, 30, 20, 18, 8, 8, 6, 1, 1, 1))

chart_frame = slide3.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(0.5), Inches(1.2), Inches(7.5), Inches(5.5),
    chart_data
)
chart = chart_frame.chart
chart.style = 2
plot = chart.plots[0]
plot.gap_width = 80
series = plot.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = ORANGE

# Data labels
series.has_data_labels = True
data_labels = series.data_labels
data_labels.font.size = Pt(10)
data_labels.font.color.rgb = TEXT_WHITE
data_labels.font.bold = True
data_labels.number_format = '0'
data_labels.show_value = True

chart.has_legend = False

# Category axis
cat_axis = chart.category_axis
cat_axis.tick_labels.font.size = Pt(10)
cat_axis.tick_labels.font.color.rgb = TEXT_GRAY

# Value axis
val_axis = chart.value_axis
val_axis.tick_labels.font.size = Pt(9)
val_axis.tick_labels.font.color.rgb = TEXT_DIM

# Right side - breakdown
add_shape(slide3, Inches(8.4), Inches(1.2), Inches(4.5), Inches(5.5), BG_CARD, RGBColor(0x1E, 0x29, 0x3B))

add_text_box(slide3, Inches(8.7), Inches(1.4), Inches(4), Inches(0.4),
             'Project Breakdown', font_size=16, color=ORANGE, bold=True)

project_details = [
    ('CBP — Builder Platform', '35+ active tickets', '5 P1 critical, 2 blocked', PURPLE),
    ('DD — Documents & Delivery', '30 active tickets', 'Jacob Powers: 10, Ya He: 8, EPD-logged', GREEN),
    ('SENG — Solution Engineering', '20 active tickets', '12 assigned to Omar Aparicio', ORANGE),
    ('CLN — Consumer Lending', '18 active tickets', '9 P1, 1 blocked (CLN-185)', BLUE),
    ('DDINDIA — D&D India', '8 active tickets', '3 P2 issues, 1 blocked (DDINDIA-555)', RGBColor(0x22, 0xD3, 0xEE)),
    ('DA / RHL / APEX / BAI / IMB', '17 active tickets', 'DA/RHL: 14, others low volume', SLATE),
]

y = Inches(2.0)
for name, count, detail, clr in project_details:
    add_shape(slide3, Inches(8.7), y + Inches(0.02), Inches(0.06), Inches(0.55), clr)
    add_text_box(slide3, Inches(8.95), y, Inches(3.9), Inches(0.3),
                 name, font_size=11, color=clr, bold=True)
    add_text_box(slide3, Inches(8.95), y + Inches(0.25), Inches(3.9), Inches(0.2),
                 count, font_size=10, color=TEXT_WHITE)
    add_text_box(slide3, Inches(8.95), y + Inches(0.42), Inches(3.9), Inches(0.2),
                 detail, font_size=9, color=TEXT_GRAY)
    y += Inches(0.72)

add_shape(slide3, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), ACCENT_RED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — Priority & Status Distribution
# ═══════════════════════════════════════════════════════════════════════════

slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, BG_DARK)
add_shape(slide4, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_RED)

add_text_box(slide4, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             'Priority & Status Distribution', font_size=28, color=RED, bold=True)

# Priority chart
priority_data = CategoryChartData()
priority_data.categories = ['P0 (Blocking)', 'P1 (Critical)', 'P2 (Major)', 'P3 (Minor)', 'P4 (Trivial)']
priority_data.add_series('Count', (2, 40, 6, 68, 15))

priority_frame = slide4.shapes.add_chart(
    XL_CHART_TYPE.PIE,
    Inches(0.5), Inches(1.3), Inches(5.8), Inches(4.5),
    priority_data
)
priority_chart = priority_frame.chart
priority_chart.has_legend = True
priority_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
priority_chart.legend.font.size = Pt(9)
priority_chart.legend.font.color.rgb = TEXT_GRAY
priority_chart.legend.include_in_layout = False

plot = priority_chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.font.size = Pt(9)
data_labels.font.color.rgb = WHITE
data_labels.font.bold = True
data_labels.number_format = '0'
data_labels.show_value = True
data_labels.show_category_name = False
data_labels.show_percentage = True

add_text_box(slide4, Inches(1.5), Inches(1.1), Inches(4), Inches(0.3),
             'Priority Distribution', font_size=14, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# Status chart
status_data = CategoryChartData()
status_data.categories = ['To Do', 'Ready for Dev', 'Dev In Progress', 'QA In Progress', 'Ready for QA', 'Blocked']
status_data.add_series('Count', (30, 47, 20, 5, 7, 5))

status_frame = slide4.shapes.add_chart(
    XL_CHART_TYPE.PIE,
    Inches(6.8), Inches(1.3), Inches(5.8), Inches(4.5),
    status_data
)
status_chart = status_frame.chart
status_chart.has_legend = True
status_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
status_chart.legend.font.size = Pt(9)
status_chart.legend.font.color.rgb = TEXT_GRAY
status_chart.legend.include_in_layout = False

plot = status_chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.font.size = Pt(9)
data_labels.font.color.rgb = WHITE
data_labels.font.bold = True
data_labels.number_format = '0'
data_labels.show_value = True
data_labels.show_percentage = True

add_text_box(slide4, Inches(7.8), Inches(1.1), Inches(4), Inches(0.3),
             'Status Distribution', font_size=14, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# Key insight box
add_shape(slide4, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.2), BG_CARD, RGBColor(0x1E, 0x29, 0x3B))
add_shape(slide4, Inches(0.5), Inches(5.9), Inches(0.06), Inches(1.2), RED)

add_text_box(slide4, Inches(0.8), Inches(6.0), Inches(3), Inches(0.3),
             'Key Insight', font_size=14, color=RED, bold=True)
add_text_box(slide4, Inches(0.8), Inches(6.35), Inches(11.8), Inches(0.7),
             '32% of tickets are P0/P1 (42 of 131). 59% remain in "To Do/Ready for Dev" with no active development. '
             'DD (30 tickets) was previously untracked. Only 7 tickets are ready for QA and 5 are blocked. Significant resolution bottleneck requires resource reallocation.',
             font_size=10, color=TEXT_WHITE)

add_shape(slide4, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), ACCENT_RED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — Top OCE Assignees
# ═══════════════════════════════════════════════════════════════════════════

slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, BG_DARK)
add_shape(slide5, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_RED)

add_text_box(slide5, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             'Top OCE Assignees', font_size=28, color=RED, bold=True)
add_text_box(slide5, Inches(0.5), Inches(0.9), Inches(12), Inches(0.3),
             'Engineers with highest active OCE ticket load', font_size=12, color=TEXT_GRAY)

# Chart
assignee_data = CategoryChartData()
assignee_data.categories = ['Omar Aparicio', 'Jacob Powers', 'Ya He', 'Gaurav Kumar',
                            'Kumar Utkarsh', 'Suyog Bhatia', 'Shreshth Malik', 'Sathish Krishnan',
                            'Bukka V Reddy', 'Adjoa Quansah', 'Will Duan', 'Amod Kumar']
assignee_data.add_series('Active Tickets', (12, 10, 8, 8, 7, 4, 4, 3, 3, 3, 3, 2))

chart_frame = slide5.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(0.5), Inches(1.5), Inches(8), Inches(5.2),
    assignee_data
)
chart = chart_frame.chart
chart.style = 2
plot = chart.plots[0]
plot.gap_width = 60
series = plot.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = BLUE

series.has_data_labels = True
data_labels = series.data_labels
data_labels.font.size = Pt(10)
data_labels.font.color.rgb = TEXT_WHITE
data_labels.font.bold = True

chart.has_legend = False
chart.category_axis.tick_labels.font.size = Pt(9)
chart.category_axis.tick_labels.font.color.rgb = TEXT_GRAY
chart.value_axis.tick_labels.font.size = Pt(9)
chart.value_axis.tick_labels.font.color.rgb = TEXT_DIM

# Right callout
add_shape(slide5, Inches(8.9), Inches(1.5), Inches(4.0), Inches(5.2), BG_CARD, RGBColor(0x1E, 0x29, 0x3B))

add_text_box(slide5, Inches(9.15), Inches(1.7), Inches(3.5), Inches(0.4),
             'Workload Concerns', font_size=16, color=ORANGE, bold=True)

concerns = [
    ('Omar Aparicio (SENG)', '12 tickets — 60% of SENG backlog. Single point of failure risk.', RED),
    ('DD Engineers', 'Jacob Powers (10), Ya He (8) — D&D on-call tickets heavily concentrated on 2 engineers.', GREEN),
    ('DDINDIA Team', 'Gaurav Kumar (8), Kumar Utkarsh (7) carry most of the India D&D workload. Utkarsh has 1 blocked P2.', RGBColor(0x22, 0xD3, 0xEE)),
    ('Unassigned P1s', 'CBP-13372, CBP-13333, CLN-260 still unassigned — need immediate owners.', RED),
    ('Recommendation', 'Redistribute OCE load. Assign P1 owners within 24h. Review DD backlog for stale items.', GREEN),
]

y = Inches(2.3)
for title, desc, clr in concerns:
    add_shape(slide5, Inches(9.15), y, Inches(0.06), Inches(0.6), clr)
    add_text_box(slide5, Inches(9.4), y, Inches(3.3), Inches(0.25),
                 title, font_size=10, color=clr, bold=True)
    add_text_box(slide5, Inches(9.4), y + Inches(0.22), Inches(3.3), Inches(0.5),
                 desc, font_size=8.5, color=TEXT_WHITE)
    y += Inches(0.75)

add_shape(slide5, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), ACCENT_RED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — P0/P1 Critical Tickets (Table)
# ═══════════════════════════════════════════════════════════════════════════

slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, BG_DARK)
add_shape(slide6, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_RED)

add_text_box(slide6, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             'P0/P1 Critical Tickets — Immediate Attention Required', font_size=24, color=RED, bold=True)

p1_tickets = [
    ['CBP-13372', 'CBP', 'CP enablement stuck nodes for many tenants', 'P1', 'To Do', 'Unassigned'],
    ['CBP-13333', 'CBP', 'Document upload screen blanks out', 'P1', 'To Do', 'Unassigned'],
    ['CBP-13234', 'CBP', 'Builder Component Enhancements Q1 2026', 'P1', 'To Do', 'Shimrit Yacobi'],
    ['CBP-11817', 'CBP', 'Step-up screen issue: page not loading', 'P1', 'Blocked', 'Kellis Liang'],
    ['CBP-7059', 'CBP', 'Lending module registry incorrect data', 'P1', 'Dev In Progress', 'Kevin Wang'],
    ['CLN-260', 'CLN', 'TeachersFCU - Vehicle details duplicated', 'P1', 'To Do', 'Unassigned'],
    ['CLN-245', 'CLN', 'TeachersFCU - Missing collateral information', 'P1', 'Dev In Progress', 'Kalpana Sharma'],
    ['CLN-237', 'CLN', 'BMO PL: Store credit pull for frozen credit', 'P1', 'To Do', 'Suyog Bhatia'],
    ['CLN-221', 'CLN', 'Align PLv6 App Data with MISMO', 'P1', 'Dev In Progress', 'Sathish Krishnan'],
    ['CLN-215', 'CLN', 'TeachersFCU - App stuck same email co-applicant', 'P1', 'Dev In Progress', 'Suyog Bhatia'],
    ['CLN-193', 'CLN', 'BMO PL - Only latest residency in DMS', 'P1', 'Dev In Progress', 'Bukka V Reddy'],
    ['CLN-190', 'CLN', 'BMO PL: Expand CB Credit Hard Pull Payload', 'P1', 'Dev In Progress', 'Bukka V Reddy'],
    ['CLN-188', 'CLN', 'BMO PL - Support eSign Followups', 'P1', 'Dev In Progress', 'Eduardo Angeles'],
    ['CLN-185', 'CLN', 'PL - Debt consolidation not appearing', 'P1', 'Blocked', 'Suyog Bhatia'],
]

headers = ['Ticket', 'Project', 'Summary', 'Priority', 'Status', 'Assignee']
col_widths = [Inches(1.1), Inches(0.7), Inches(5.4), Inches(0.8), Inches(1.4), Inches(1.8)]

add_table_slide(slide6, headers, p1_tickets,
                Inches(0.5), Inches(1.2), Inches(11.2), col_widths)

add_shape(slide6, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), ACCENT_RED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — EPD-Logged Issues by Project
# ═══════════════════════════════════════════════════════════════════════════

slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, BG_DARK)
add_shape(slide7, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_RED)

add_text_box(slide7, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             'EPD-Logged Production Issues — By Project', font_size=28, color=RED, bold=True)
add_text_box(slide7, Inches(0.5), Inches(0.9), Inches(12), Inches(0.3),
             '47 EPD-logged issues across engineering teams requiring resolution', font_size=12, color=TEXT_GRAY)

# Project breakdown cards
projects = [
    ('DD — Documents & Delivery', '30 active', GREEN, [
        'DD-3514: QA Updates to eSign Follow-ups in CP (EPD-logged)',
        'DD-3522: Manual Follow-ups not working CP V5 PL/CC (EPD)',
        'DD-4079: Inconsistent Scheduling Task — uwcreditunion',
        'DD-3803: Cannot access borrower Identity Docs in Signing',
        'DD-3726: Disclosures messages table growing rapidly',
        '+ 25 more (Jacob Powers: 10, Ya He: 8, Adjoa: 3)',
    ]),
    ('DDINDIA — D&D India', '8 active', RGBColor(0x22, 0xD3, 0xEE), [
        'DDINDIA-548: Truist MOA not recognizing ICD (P2, In Dev)',
        'DDINDIA-555: 4506C doc prompting wrong borrower (P2, Blocked)',
        'DDINDIA-302: Penfed — unhandled server error (P2, In Dev)',
        'DDINDIA-558: LOE Large Deposits doc issue — BMO',
        'DDINDIA-93: W2 DocAI not verifying on first pass (EPD)',
        '+ 3 more (Gaurav Kumar, Kumar Utkarsh, Monilkumar)',
    ]),
    ('CLN — Consumer Lending', '18 active', BLUE, [
        'CLN-260: TFCU - Vehicle details duplicated on review (P1)',
        'CLN-248: PLv6 areBorrowersMarried checks primary only (P1)',
        'CLN-245: TFCU - Missing collateral info (P1, In Dev)',
        'CLN-237: BMO PL - Store credit pull for frozen credit (P1)',
        'CLN-185: PL - Debt consolidation not appearing (P1, Blocked)',
        '+ 8 more active bug tickets',
    ]),
    ('CBP — Builder Platform', '35+ active', PURPLE, [
        'CBP-13372: CP enablement stuck nodes - many tenants (P1)',
        'CBP-13333: Document upload screen blanks out (P1)',
        'CBP-13387: Subworkflow TextInput in loop incorrect value',
        'CBP-13386: Subworkflows text input not refreshed',
        'CBP-11817: Step-up screen issue page not loading (P1, Blocked)',
        '+ 25 more active platform bugs',
    ]),
]

# 2x2 grid
positions = [
    (Inches(0.4), Inches(1.5)),   # top-left
    (Inches(6.7), Inches(1.5)),   # top-right
    (Inches(0.4), Inches(4.5)),   # bottom-left
    (Inches(6.7), Inches(4.5)),   # bottom-right
]

for i, (name, count, clr, items) in enumerate(projects):
    x, y = positions[i]
    card_w, card_h = Inches(6.0), Inches(2.8)
    add_shape(slide7, x, y, card_w, card_h, BG_CARD, RGBColor(0x1E, 0x29, 0x3B))
    add_shape(slide7, x, y, Inches(0.06), card_h, clr)

    add_text_box(slide7, x + Inches(0.2), y + Inches(0.1), Inches(4.5), Inches(0.3),
                 name, font_size=12, color=clr, bold=True)
    add_text_box(slide7, x + Inches(4.8), y + Inches(0.1), Inches(1.0), Inches(0.3),
                 count, font_size=10, color=RED, bold=True, alignment=PP_ALIGN.RIGHT)

    item_y = y + Inches(0.5)
    for item in items:
        add_text_box(slide7, x + Inches(0.25), item_y, Inches(5.5), Inches(0.25),
                     f'• {item}', font_size=7.5, color=TEXT_WHITE)
        item_y += Inches(0.28)

add_shape(slide7, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), ACCENT_RED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — Ticket Categories & Trends
# ═══════════════════════════════════════════════════════════════════════════

slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, BG_DARK)
add_shape(slide8, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_RED)

add_text_box(slide8, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             'Ticket Categories & Issue Classification', font_size=28, color=RED, bold=True)

# Category chart
cat_data = CategoryChartData()
cat_data.categories = ['EPD-Logged', 'P1 Bugs', 'DD On-Call', 'Enterprise UI', 'DDINDIA On-Call', 'OCE On-Call', 'DA V12', 'Config/Tmpl', 'Customer-QA', 'Subworkflow']
cat_data.add_series('Count', (52, 42, 30, 12, 8, 4, 6, 5, 2, 3))

cat_frame = slide8.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.5), Inches(1.3), Inches(7), Inches(5.2),
    cat_data
)
cat_chart = cat_frame.chart
cat_chart.style = 2
plot = cat_chart.plots[0]
plot.gap_width = 80
series = plot.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = PURPLE

series.has_data_labels = True
dl = series.data_labels
dl.font.size = Pt(9)
dl.font.color.rgb = TEXT_WHITE
dl.font.bold = True
dl.label_position = XL_LABEL_POSITION.OUTSIDE_END

cat_chart.has_legend = False
cat_chart.category_axis.tick_labels.font.size = Pt(8)
cat_chart.category_axis.tick_labels.font.color.rgb = TEXT_GRAY
cat_chart.value_axis.tick_labels.font.size = Pt(8)
cat_chart.value_axis.tick_labels.font.color.rgb = TEXT_DIM

# Key observations
add_shape(slide8, Inches(8), Inches(1.3), Inches(4.9), Inches(5.2), BG_CARD, RGBColor(0x1E, 0x29, 0x3B))

add_text_box(slide8, Inches(8.3), Inches(1.5), Inches(4.3), Inches(0.4),
             'Category Analysis', font_size=16, color=ORANGE, bold=True)

categories_text = [
    ('EPD-Logged (52)', 'Largest category. Production defects logged via EPD across all projects including DD and DDINDIA.', RED),
    ('P1 Critical Bugs (42)', 'High priority items in CBP and CLN that directly impact customers and product stability.', ORANGE),
    ('DD On-Call (30)', 'Documents & Delivery on-call form tickets — disclosures, eSign, docgen issues. Jacob Powers & Ya He carry most load.', GREEN),
    ('DDINDIA On-Call (8)', 'India D&D team — 3 P2 issues (Truist MOA, 4506C doc, Penfed disclosures). 1 blocked.', RGBColor(0x22, 0xD3, 0xEE)),
    ('Enterprise UI (12)', 'SENG — Omar Aparicio. Enterprise UI issues from withdrawn apps, login, duplicate elements.', BLUE),
]

y = Inches(2.1)
for title, desc, clr in categories_text:
    add_shape(slide8, Inches(8.3), y, Inches(0.06), Inches(0.58), clr)
    add_text_box(slide8, Inches(8.55), y, Inches(4.1), Inches(0.22),
                 title, font_size=9.5, color=clr, bold=True)
    add_text_box(slide8, Inches(8.55), y + Inches(0.2), Inches(4.1), Inches(0.45),
                 desc, font_size=8, color=TEXT_WHITE)
    y += Inches(0.66)

add_shape(slide8, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), ACCENT_RED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — Recommendations & Action Items
# ═══════════════════════════════════════════════════════════════════════════

slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9, BG_DARK)
add_shape(slide9, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_RED)

add_text_box(slide9, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             'Recommendations & Action Items', font_size=28, color=RED, bold=True)

# Three columns of recommendations
recommendations = [
    {
        'title': 'Immediate (This Week)',
        'color': RED,
        'items': [
            'Assign owners to unassigned P1s (CBP-13372, CBP-13333, CLN-260, DA-1612)',
            'Unblock DDINDIA-555 (4506C doc issue for Wyndham) — P2 customer impact',
            'Unblock CBP-11817, CLN-185, RHL-1659 — dependency resolution needed',
            'Triage DD backlog: 30 on-call tickets, many >3 months old — close stale items',
            'Triage CBP tickets: prioritize CP enablement (CBP-13372) for multiple tenants',
        ]
    },
    {
        'title': 'Short-Term (2 Weeks)',
        'color': ORANGE,
        'items': [
            'Redistribute DD load — Jacob Powers (10) and Ya He (8) carry 60% of DD backlog',
            'Redistribute SENG load — Omar Aparicio at 12 tickets is unsustainable',
            'Accelerate CLN P1 resolution — 9 P1s need sprint allocation',
            'DDINDIA: Resolve 3 P2 customer issues (Truist, Penfed, Wyndham)',
            'Move 7 "Ready for QA" tickets through QA cycle to reduce backlog',
        ]
    },
    {
        'title': 'Process Improvements',
        'color': GREEN,
        'items': [
            'Set SLA targets: P1 = 5 business days, P0 = 24 hours',
            'Include DD & DDINDIA in regular OCE tracking — previously unmonitored',
            'Create cross-project OCE standup (15 min daily) for P0/P1 tracking',
            'Archive stale DD tickets (>6 months old, P4) to reduce noise',
            'Monthly OCE retrospective to identify recurring issue patterns',
        ]
    }
]

x_positions = [Inches(0.4), Inches(4.55), Inches(8.7)]
for i, rec in enumerate(recommendations):
    x = x_positions[i]
    card_w = Inches(3.95)
    card_h = Inches(5.3)
    add_shape(slide9, x, Inches(1.3), card_w, card_h, BG_CARD, RGBColor(0x1E, 0x29, 0x3B))
    add_shape(slide9, x, Inches(1.3), card_w, Inches(0.06), rec['color'])

    add_text_box(slide9, x + Inches(0.2), Inches(1.5), Inches(3.5), Inches(0.35),
                 rec['title'], font_size=14, color=rec['color'], bold=True)

    y = Inches(2.0)
    for item in rec['items']:
        # Number circle
        idx = rec['items'].index(item) + 1
        add_text_box(slide9, x + Inches(0.15), y, Inches(3.6), Inches(0.5),
                     f'{idx}.  {item}', font_size=8.5, color=TEXT_WHITE)
        y += Inches(0.6)

# Summary bar at bottom
add_shape(slide9, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.5), BG_CARD, RGBColor(0x1E, 0x29, 0x3B))
add_shape(slide9, Inches(0.4), Inches(6.75), Inches(0.06), Inches(0.5), AMBER)
add_text_box(slide9, Inches(0.7), Inches(6.8), Inches(12), Inches(0.4),
             'Goal: Reduce active OCE tickets from 131 to <70 within 4 weeks by prioritizing P0/P1 resolution, triaging DD backlog, unblocking dependencies, and redistributing workload.',
             font_size=10, color=TEXT_WHITE)

add_shape(slide9, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), ACCENT_RED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — Thank You / Q&A
# ═══════════════════════════════════════════════════════════════════════════

slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10, BG_DARK)
add_shape(slide10, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_RED)

add_text_box(slide10, Inches(1), Inches(2.2), Inches(11.3), Inches(1),
             'Questions & Discussion', font_size=40, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide10, Inches(1), Inches(3.3), Inches(11.3), Inches(0.5),
             'OCE Active Tickets  |  Blend Engineering', font_size=16, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Stats summary
summary_items = [
    ('131 Active Tickets', RED),
    ('42 P0/P1 Critical', ORANGE),
    ('52 EPD-Logged', AMBER),
    ('10 Projects', BLUE),
    ('5 Blocked', PURPLE),
]

total_w = len(summary_items) * 2.2
start_x = (13.333 - total_w) / 2
for i, (text, clr) in enumerate(summary_items):
    x = Inches(start_x + i * 2.2)
    add_shape(slide10, x, Inches(4.3), Inches(2.0), Inches(0.5), BG_CARD, RGBColor(0x1E, 0x29, 0x3B))
    add_text_box(slide10, x + Inches(0.05), Inches(4.35), Inches(1.9), Inches(0.4),
                 text, font_size=10, color=clr, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide10, Inches(1), Inches(5.5), Inches(11.3), Inches(0.4),
             f'Data sourced from JIRA (blendlabs.atlassian.net)  ·  Generated: {now}', font_size=11, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_shape(slide10, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), ACCENT_RED)


# ═══════════════════════════════════════════════════════════════════════════
#  Save
# ═══════════════════════════════════════════════════════════════════════════

output_path = '/Users/vinay-prasadg/Documents/Production Defects/OCE-Active-Tickets-Report.pptx'
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Slides: {len(prs.slides)}')
